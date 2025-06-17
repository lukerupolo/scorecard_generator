import streamlit as st
import requests
import pandas as pd
import numpy as np
from datetime import datetime, timedelta
from io import BytesIO

# --- PowerPoint and Styling Imports ---
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt

# --- Excel Styling Imports ---
from openpyxl.styles import Font, PatternFill, Alignment

# ================================================================================
# TODO: PASTE YOUR EXTRACTED STYLE INFORMATION HERE
# ================================================================================
# Use the RGB values and font names you got from the extract_style.py script
BRAND_STYLE = {
    "colors": {
        "primary":      RGBColor(79, 129, 189),   # Example: Your 'Accent 1' color
        "accent":       RGBColor(192, 80, 77),    # Example: Your 'Accent 2' color
        "text_dark":    RGBColor(38, 38, 38),     # Example: Your 'Dark 1' text color
        "text_light":   RGBColor(255, 255, 255),  # Example: Your 'Light 1' text color
        "background_alt": RGBColor(242, 242, 242),  # A light grey for table rows
    },
    "fonts": {
        "heading": "Calibri",  # Example: The 'Major Font' you extracted
        "body": "Calibri",     # Example: The 'Minor Font' you extracted
    },
    "font_sizes": {
        "title": Pt(36),
        "subtitle": Pt(24),
        "body": Pt(12),
        "table_header": Pt(11),
        "table_body": Pt(10),
    }
}
# ================================================================================


st.set_page_config(page_title="Event Marketing Scorecard", layout="wide")
DEBUG = st.sidebar.checkbox("🔍 Show LevelUp raw data")


# ────────────────────────────────────────────────────────────────────────────────
# 1) Helper functions for LevelUp API integration and Social Mentions (Onclusive)
# ────────────────────────────────────────────────────────────────────────────────
def setup_levelup_headers(api_key: str) -> dict:
    return {
        "accept": "application/json",
        "X-API-KEY": api_key
    }

def generate_levelup_metrics_for_event(event: dict, headers: dict) -> dict:
    """
    Fetch 7-day baseline & 7-day actual Video Views & Stream Hours
    from LevelUp, using the same endpoints as your ETL script.
    """
    base_url   = "https://www.levelup-analytics.com/api/client/v1"  # ← your real base URL
    brand_id   = event["brandId"]
    region     = event["region"]
    event_date = event["date"].date()

    windows = {
        "baseline": (
            (event_date - timedelta(days=7)).strftime("%Y-%m-%d"),
            (event_date - timedelta(days=1)).strftime("%Y-%m-%d"),
        ),
        "actual": (
            event_date.strftime("%Y-%m-%d"),
            (event_date + timedelta(days=6)).strftime("%Y-%m-%d"),
        ),
    }

    def fetch_metrics(metric: str, start: str, end: str, period: str) -> pd.DataFrame:
        url    = f"{base_url}/{metric}/statsEvolution/brand/{brand_id}"
        params = {"from": start, "to": end, "brandid": brand_id, "region": region}
        try:
            r = requests.get(url, headers=headers, params=params)
            r.raise_for_status()
            df = pd.DataFrame(r.json().get("data", []))
            if not df.empty:
                df["period"] = period
            return df
        except Exception as e:
            print(f"⚠️ Error fetching {metric} ({period}): {e}")
            return pd.DataFrame()

    videos = [fetch_metrics("videos", s, e, p) for p, (s, e) in windows.items()]
    streams = [fetch_metrics("streams", s, e, p) for p, (s, e) in windows.items()]

    return {"videos": pd.concat(videos, ignore_index=True),
            "streams": pd.concat(streams, ignore_index=True)}

def compute_three_month_average(
    headers: dict,
    brand_id: int,
    region: str,
    event_date: datetime,
    metric: str,
) -> float:
    """
    Compute the 12-week average for 'videos' or 'streams'
    over the 90 days before the event.
    """
    base_url   = "https://www.levelup-analytics.com/api/client/v1"  # ← and here
    end_date   = (event_date - timedelta(days=1)).strftime("%Y-%m-%d")
    start_date = (event_date - timedelta(days=90)).strftime("%Y-%m-%d")

    url    = f"{base_url}/{metric}/statsEvolution/brand/{brand_id}"
    params = {"from": start_date, "to": end_date, "brandid": brand_id, "region": region}

    try:
        r = requests.get(url, headers=headers, params=params)
        r.raise_for_status()
        df = pd.DataFrame(r.json().get("data", []))
        if df.empty:
            return 0.0
        if metric == "videos":
            return df["views"].sum() / 12
        col = "hoursWatched" if "hoursWatched" in df.columns else "watchTime"
        return df[col].sum() / 12
    except Exception as e:
        print(f"⚠️ Error computing 3-month avg for {metric}: {e}")
        return 0.0

# ────────────────────────────────────────────────────────────────────────────────
# 2) Streamlit app configuration and sidebar
# ────────────────────────────────────────────────────────────────────────────────
# Predefined game-to-ID mapping and regions list for dropdowns
game_options = {
    "EA Sports FC25": 3136,
    "FIFA 25": 3140,
    "Madden NFL 25": 3150,
    "NHL 25": 3160,
}
# Add ‘Other’ to allow custom region codes
region_options = ["US", "GB", "AU", "CA", "FR", "DE", "JP", "KR", "TH", "Other"]

# ────────────────────────────────────────────────────────────────────────────────
# 3) Sidebar: Event Configuration with dropdowns
# ────────────────────────────────────────────────────────────────────────────────
st.sidebar.markdown("## 📅 Event Configuration")
st.sidebar.markdown("Set up one or more events and their details below.")

n_events = st.sidebar.number_input(
    "Number of events", min_value=1, max_value=10, value=1, step=1
)
events: list[dict] = []
for i in range(n_events):
    st.sidebar.markdown(f"### Event {i+1}")
    
    # Dropdown for game selection
    selected_game = st.sidebar.selectbox(
        f"🎮 Select Game (Event {i+1})",
        options=list(game_options.keys()),
        key=f"game_{i}"
    )
    brand_name = selected_game
    brand_id = game_options[selected_game]

    # Event label (defaults to game name)
    name = st.sidebar.text_input(
        f"🔤 Event Label (Event {i+1})",
        key=f"name_{i}",
        value=brand_name
    )
    # Date picker
    date = st.sidebar.date_input(
        f"📅 Date (Event {i+1})", key=f"date_{i}"
    )

    # Dropdown for region code with custom entry
    selected_region = st.sidebar.selectbox(
        f"🌐 Select Region (Event {i+1})",
        options=region_options,
        key=f"region_select_{i}"
    )
    if selected_region == "Other":
        region = st.sidebar.text_input(
            f"Enter custom region code (Event {i+1})",
            key=f"region_custom_{i}"
        ) or ""
    else:
        region = selected_region

    events.append({
        "name": name,
        "date": datetime.combine(date, datetime.min.time()),
        "brandId": int(brand_id),
        "brandName": brand_name,
        "region": region,
    })

# ────────────────────────────────────────────────────────────────────────────────
# 4) Sidebar: Metric Selection with custom metric input
# ────────────────────────────────────────────────────────────────────────────────
st.sidebar.markdown("## 🎛️ Metric Selection")
st.sidebar.markdown("Choose one or more predefined metrics or add your own custom metric.")

predefined_metrics = [
    "Video Views (VOD)", "Hours Watched (Streams)", "Social Mentions", "Sessions",
    "DAU", "Revenue", "Installs", "Retention", "Watch Time", "ARPU", "Conversions",
    "Search Index", "PCCV", "AMA", "Stream Views", "UGC Views",
    "Social Impressions (FC Owned Channels)", "Social Conversation Volume", "Social Sentiment",
]
metrics = st.sidebar.multiselect("Select metrics:", options=predefined_metrics, default=[], key="metrics_multiselect")
custom_metric = st.sidebar.text_input("✍️ Add Custom Metric", placeholder="Type a custom metric and press Enter", key="custom_metric_input")
if custom_metric:
    metrics.append(custom_metric)

# ────────────────────────────────────────────────────────────────────────────────
# 5 & 6) Authentication Sections
# ────────────────────────────────────────────────────────────────────────────────
# ... (Your existing Onclusive and LevelUp authentication code sections go here) ...
# This has been omitted for brevity but should be included in your final script.

# Initialize session_state flags
if "scorecard_ready" not in st.session_state:
    st.session_state["scorecard_ready"] = False
if "sheets_dict" not in st.session_state:
    st.session_state["sheets_dict"] = {}
if "presentation_buffer" not in st.session_state:
    st.session_state["presentation_buffer"] = None

# ────────────────────────────────────────────────────────────────────────────────
# 7) Main: Generate Scorecard
# ────────────────────────────────────────────────────────────────────────────────
if st.button("✅ Generate Scorecard"):
    # ... (Your existing scorecard generation logic goes here) ...
    # This has been omitted for brevity but should be included in your final script.
    # The important part is that it populates st.session_state["sheets_dict"]
    # and sets st.session_state["scorecard_ready"] = True
    # For demonstration, we'll create dummy data if the original logic is not present
    if not st.session_state["sheets_dict"]:
        st.warning("Running with dummy data for demonstration.")
        dummy_df = pd.DataFrame({
            'Metric': ['Video Views (VOD)', 'Hours Watched (Streams)', 'Social Mentions'],
            'Baseline': [1000, 500, 200],
            'Actual': [1500, 750, 300]
        })
        st.session_state["sheets_dict"] = {"Dummy Event": dummy_df}
        st.session_state["scorecard_ready"] = True

# ────────────────────────────────────────────────────────────────────────────────
# 8) Download Excel
# ────────────────────────────────────────────────────────────────────────────────
def style_excel_sheet(writer, sheet_name, df):
    """Applies BRAND_STYLE to an Excel sheet using openpyxl."""
    workbook = writer.book
    worksheet = writer.sheets[sheet_name]

    # Convert RGBColor to hex string for openpyxl
    primary_hex = f'{BRAND_STYLE["colors"]["primary"]}'
    alt_bg_hex = f'{BRAND_STYLE["colors"]["background_alt"]}'
    text_light_hex = f'{BRAND_STYLE["colors"]["text_light"]}'

    header_font = Font(name=BRAND_STYLE["fonts"]["heading"], size=11, bold=True, color=text_light_hex)
    header_fill = PatternFill(start_color=primary_hex, end_color=primary_hex, fill_type="solid")
    alt_row_fill = PatternFill(start_color=alt_bg_hex, end_color=alt_bg_hex, fill_type="solid")

    # Style header
    for cell in worksheet[1]:
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center', vertical='center')

    # Style data rows
    for row_idx in range(2, worksheet.max_row + 1):
        if row_idx % 2 == 0:
            for cell in worksheet[row_idx]:
                cell.fill = alt_row_fill

if st.session_state.get("scorecard_ready", False):
    sheets_dict = st.session_state["sheets_dict"]
    buffer = BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        for sheet_name, df_sheet in sheets_dict.items():
            safe_name = sheet_name[:31]
            df_sheet.to_excel(writer, sheet_name=safe_name, index=False)
            # Apply the styling function to the newly created sheet
            style_excel_sheet(writer, safe_name, df_sheet)
            
    buffer.seek(0)
    st.download_button(label="📥 Download Full Scorecard Workbook", data=buffer, file_name="event_marketing_scorecard.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

# ────────────────────────────────────────────────────────────────────────────────
# 9) PowerPoint Generation Functions
# ────────────────────────────────────────────────────────────────────────────────
def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Assumes a standard Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_timeline_slide(prs, timeline_moments):
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_shape.text = "Scorecard Moments Timeline"
    p = title_shape.text_frame.paragraphs[0]
    p.font.name = BRAND_STYLE["fonts"]["heading"]
    p.font.size = BRAND_STYLE["font_sizes"]["title"]
    p.font.color.rgb = BRAND_STYLE["colors"]["text_dark"]

    if not timeline_moments: return

    # Create timeline plot with brand colors
    fig, ax = plt.subplots(figsize=(10, 2))
    ax.axhline(0, color=f'#{BRAND_STYLE["colors"]["text_dark"]}', xmin=0.05, xmax=0.95, zorder=1)
    for i, moment in enumerate(timeline_moments):
        ax.plot(i + 1, 0, 'o', markersize=20, color=f'#{BRAND_STYLE["colors"]["primary"]}', zorder=2)
        ax.text(i + 1, -0.4, moment, ha='center', fontsize=12, fontname=BRAND_STYLE["fonts"]["body"])
    ax.axis('off')

    plot_stream = BytesIO()
    plt.savefig(plot_stream, format='png', bbox_inches='tight', pad_inches=0, transparent=True)
    plt.close(fig)
    plot_stream.seek(0)
    slide.shapes.add_picture(plot_stream, Inches(0.5), Inches(1.5), width=Inches(9))

def add_moment_title_slide(prs, title_text):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text

def apply_table_style_pptx(table):
    """Applies branding from BRAND_STYLE to a python-pptx table."""
    # Style header row
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_STYLE["colors"]["primary"]
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = BRAND_STYLE["colors"]["text_light"]
        p.font.name = BRAND_STYLE["fonts"]["heading"]
        p.font.size = BRAND_STYLE["font_sizes"]["table_header"]
        p.alignment = PP_ALIGN.CENTER

    # Style data rows
    for row_idx, row in enumerate(table.rows[1:], start=1):
        fill = PatternFill(start_color=f'{BRAND_STYLE["colors"]["background_alt"]}', end_color=f'{BRAND_STYLE["colors"]["background_alt"]}', fill_type="solid")
        if row_idx % 2 != 0: # Apply alternating color
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRAND_STYLE["colors"]["background_alt"]

        for cell in row.cells:
            p = cell.text_frame.paragraphs[0]
            p.font.name = BRAND_STYLE["fonts"]["body"]
            p.font.size = BRAND_STYLE["font_sizes"]["table_body"]
            p.font.color.rgb = BRAND_STYLE["colors"]["text_dark"]


def add_df_to_slide(prs, df, slide_title):
    slide_layout = prs.slide_layouts[5] # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_shape.text = slide_title
    p = title_shape.text_frame.paragraphs[0]
    p.font.name = BRAND_STYLE["fonts"]["heading"]
    p.font.size = Pt(28)
    p.font.color.rgb = BRAND_STYLE["colors"]["text_dark"]

    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.8)).table
    table.columns[0].width = Inches(2.5)
    for i in range(1, cols):
        table.columns[i].width = Inches(1.5)

    # Add header row from DataFrame columns
    for i, col_name in enumerate(df.columns):
        table.cell(0, i).text = col_name

    # Add data rows from DataFrame
    for r in range(rows):
        for c in range(cols):
            cell_value = df.iloc[r, c]
            table.cell(r + 1, c).text = str(cell_value)
    
    # Apply the brand styling to the table
    apply_table_style_pptx(table)

# ────────────────────────────────────────────────────────────────────────────────
# 10) Main: Generate PowerPoint Presentation
# ────────────────────────────────────────────────────────────────────────────────
if st.session_state.get("scorecard_ready", False):
    if st.button("📊 Create a Game Scorecard Presentation"):
        st.session_state['show_ppt_creator'] = True

if st.session_state.get('show_ppt_creator', False):
    with st.form("ppt_form"):
        st.subheader("Presentation Details")
        ppt_title = st.text_input("Presentation Title", "Game Scorecard")
        ppt_subtitle = st.text_input("Presentation Subtitle", "A detailed analysis")
        
        scorecard_moments = st.multiselect(
            "Select Scorecard Moments for Timeline",
            options=["Pre-Reveal", "Reveal", "Pre-Order", "Launch", "Post-Launch"],
            default=["Pre-Reveal", "Reveal", "Launch"]
        )
        submitted = st.form_submit_button("Generate Presentation")

        if submitted:
            if not st.session_state["sheets_dict"]:
                st.error("No scorecard data available to create a presentation.")
            else:
                # =================================================================
                # TODO: REPLACE WITH YOUR PRESENTATION FILENAME
                # =================================================================
                try:
                    prs = Presentation("my_brand_presentation.pptx")
                except Exception as e:
                    st.warning(f"Could not load template 'my_brand_presentation.pptx'. Using a default template. Error: {e}")
                    prs = Presentation()
                # =================================================================

                # Slide 1: Title Slide
                add_title_slide(prs, ppt_title, ppt_subtitle)
                # Slide 2: Timeline
                add_timeline_slide(prs, scorecard_moments)
                
                # Subsequent slides for each scorecard moment
                for moment in scorecard_moments:
                    add_moment_title_slide(prs, f"Scorecard: {moment}")
                    # Use the actual scorecard data from the session state
                    for sheet_name, scorecard_df in st.session_state["sheets_dict"].items():
                         if sheet_name.lower() != "benchmark":
                            add_df_to_slide(prs, scorecard_df, f"{moment} Metrics: {sheet_name}")

                # Add Benchmark slide if it exists
                if "Benchmark" in st.session_state["sheets_dict"]:
                     add_moment_title_slide(prs, "Proposed Benchmarks")
                     benchmark_df = st.session_state["sheets_dict"]["Benchmark"]
                     add_df_to_slide(prs, benchmark_df, "Benchmark Summary")


                # Save presentation to a BytesIO object
                ppt_buffer = BytesIO()
                prs.save(ppt_buffer)
                ppt_buffer.seek(0)
                st.session_state["presentation_buffer"] = ppt_buffer
                st.success("Presentation generated successfully!")

if st.session_state["presentation_buffer"]:
    st.download_button(
        label="📥 Download PowerPoint Presentation",
        data=st.session_state["presentation_buffer"],
        file_name="game_scorecard_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
