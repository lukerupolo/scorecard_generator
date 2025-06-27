import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import io
import copy
import uuid
import openai
import json

# --- Core PowerPoint Functions ---

def deep_copy_slide_content(dest_slide, src_slide):
    """
    Performs a stable deep copy of all shapes from a source slide to a
    destination slide, handling different shape types robustly.
    This approach aims to minimize repair issues by using python-pptx's API
    for common shape types, especially images and text.
    """
    # Clear all shapes from the destination slide first to prepare it.
    # This loop safely removes shapes by iterating on a copy of the shapes list.
    for shape in list(dest_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    for shape in src_slide.shapes:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # For pictures, extract the image data and re-add it using python-pptx's API.
            # This is CRUCIAL for avoiding repair issues with images and ensuring proper embedding.
            try:
                image_bytes = shape.image.blob
                dest_slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
            except Exception as e:
                # Log if an image cannot be copied, but continue with other shapes
                print(f"Warning: Could not copy picture from source slide. Error: {e}")
                # Fallback: if picture has a placeholder, try to copy its XML
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    new_el = copy.deepcopy(shape.element)
                    dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                
        elif shape.has_text_frame:
            # Create a new text box on the destination slide with the same dimensions
            new_shape = dest_slide.shapes.add_textbox(left, top, width, height)
            new_text_frame = new_shape.text_frame
            new_text_frame.clear() # Clear existing paragraphs to ensure a clean copy

            # Copy text and formatting paragraph by paragraph, run by run
            for paragraph in shape.text_frame.paragraphs:
                new_paragraph = new_text_frame.add_paragraph()
                # Copy paragraph properties (e.g., alignment, indentation)
                new_paragraph.alignment = paragraph.alignment
                if hasattr(paragraph, 'level'): # Bullet level
                    new_paragraph.level = paragraph.level
                
                # Copy runs with their font properties
                for run in paragraph.runs:
                    new_run = new_paragraph.add_run()
                    new_run.text = run.text
                    
                    # Copy essential font properties (bold, italic, underline, size)
                    new_run.font.bold = run.font.bold
                    new_run.font.italic = run.font.italic
                    new_run.font.underline = run.font.underline
                    if run.font.size: # Only copy if size is explicitly defined
                        new_run.font.size = run.font.size
                    
                    # Copy font color if it's a solid fill RGB color
                    if run.font.fill.type == 1: # MSO_FILL_TYPE.SOLID
                        new_run.font.fill.solid()
                        try:
                            # Ensure color is an RGBColor object for direct assignment
                            if isinstance(run.font.fill.fore_color.rgb, RGBColor):
                                new_run.font.fill.fore_color.rgb = run.font.fill.fore_color.rgb
                            else: 
                                # Attempt to convert to RGBColor if not already
                                # This handles cases where color might be a theme color or other type
                                rgb_tuple = run.font.fill.fore_color.rgb # Assuming it might be a tuple (R, G, B)
                                new_run.font.fill.fore_color.rgb = RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])
                        except Exception as color_e:
                            print(f"Warning: Could not copy font color. Error: {color_e}")
                            pass # If color conversion fails, skip copying the color

            # Copy text frame properties (word wrap, margins)
            new_text_frame.word_wrap = shape.text_frame.word_wrap
            new_text_frame.margin_left = shape.text_frame.margin_left
            new_text_frame.margin_right = shape.text_frame.margin_right
            new_text_frame.margin_top = shape.text_frame.margin_top
            new_text_frame.margin_bottom = shape.text_frame.margin_bottom

        else:
            # For other shapes (e.g., simple geometric shapes, lines, groups, tables, charts),
            # fall back to deep copying the raw XML element.
            # This is less robust than using specific python-pptx add_* methods but necessary
            # for types not directly supported by add_*.
            # For complex custom shapes, this might still lead to minor issues,
            # but is the best general approach without parsing deeper XML.
            new_el = copy.deepcopy(shape.element)
            dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

def find_slide_by_ai(api_key, prs, slide_type_prompt, deck_name):
    """
    Uses OpenAI to intelligently find the best matching slide and get a justification.
    Returns a dictionary with the slide object, its index, and the AI's justification.
    """
    if not slide_type_prompt: return {"slide": None, "index": -1, "justification": "No keyword provided."}
    
    # Check if API key is provided and valid
    if not api_key:
        return {"slide": None, "index": -1, "justification": "OpenAI API Key is missing."}

    client = openai.OpenAI(api_key=api_key)
    
    slides_content = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        # Concatenate all text from the slide, limiting to first 1000 characters to save tokens
        slides_content.append({"slide_index": i, "text": " ".join(slide_text)[:1000]})

    system_prompt = f"""
    You are an expert presentation analyst. Your task is to find the best slide in a presentation that matches a user's description.
    The user is looking for a slide representing: '{slide_type_prompt}'.
    Analyze the text of each slide to understand its purpose. A "Timeline" slide VISUALLY represents a schedule with dates, quarters, or sequential phases (Phase 1, Phase 2); it is NOT just a list in a table of contents. An "Objectives" slide will contain goal-oriented language. You must prioritize actual content slides over simple divider or table of contents pages.
    You MUST return a JSON object with two keys: 'best_match_index' (an integer, or -1 if no match) and 'justification' (a brief, one-sentence justification for your choice).
    """
    full_user_prompt = f"Find the best slide for '{slide_type_prompt}' in the '{deck_name}' deck with the following contents:\n{json.dumps(slides_content, indent=2)}"

    try:
        response = client.chat.completions.create(
            model="gpt-4-turbo", 
            messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": full_user_prompt}],
            response_format={"type": "json_object"}
        )
        result = json.loads(response.choices[0].message.content)
        best_index = result.get("best_match_index", -1)
        justification = result.get("justification", "No justification provided.")
        
        # Validate the AI's chosen index
        if best_index != -1 and best_index < len(prs.slides):
            return {"slide": prs.slides[best_index], "index": best_index, "justification": justification}
        else:
            return {"slide": None, "index": -1, "justification": "AI could not find a suitable slide or returned an invalid index."}
    except openai.APIError as e:
        return {"slide": None, "index": -1, "justification": f"OpenAI API Error: {e}"}
    except json.JSONDecodeError as e:
        return {"slide": None, "index": -1, "justification": f"AI response was not valid JSON: {e}"}
    except Exception as e:
        return {"slide": None, "index": -1, "justification": f"An unexpected error occurred during AI analysis: {e}"}

def get_slide_content(slide):
    """Extracts title and body text from a slide."""
    if not slide: return {"title": "", "body": ""}
    
    # Sort text shapes by their top position to infer order (title usually highest)
    text_shapes = sorted([s for s in slide.shapes if s.has_text_frame and s.text.strip()], key=lambda s: s.top)
    
    title = ""
    body = ""
    
    if text_shapes:
        # Heuristic for title: often the first (top-most) text shape.
        # Could be improved by checking placeholder type (e.g., MSO_PLACEHOLDER_TYPE.TITLE)
        title = text_shapes[0].text.strip()
        body = "\n".join(s.text.strip() for s in text_shapes[1:])
        
    return {"title": title, "body": body}

def populate_slide(slide, content):
    """
    Populates a slide's placeholders or main text boxes with new content.
    It clears the existing content and adds new runs, aiming to use existing
    placeholders without forcing bold.
    """
    title_populated, body_populated = False, False
    
    # Iterate through shapes to find suitable places for title and body
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        # Check if it's a title placeholder (type 1, 2, or object type 8 which can be title)
        # Or if it's a top-positioned shape likely to be a title
        is_title_placeholder = (
            hasattr(shape, 'is_placeholder') and shape.is_placeholder and 
            shape.placeholder_format.type in (1, 2, 8) # TITLE, CENTER_TITLE, OBJECT
        )
        is_top_text_box = (shape.top < Pt(150)) # Heuristic: within 1.5 inches from top

        if not title_populated and (is_title_placeholder or is_top_text_box):
            tf = shape.text_frame
            tf.clear() # Clear existing content
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = content.get("title", "")
            # No longer forcing bold here. The template's default formatting will apply.
            title_populated = True
            
        # Check for body placeholders (type 3, 4, 8, 14) or large text boxes with dummy text
        is_body_placeholder = (
            hasattr(shape, 'is_placeholder') and shape.is_placeholder and 
            shape.placeholder_format.type in (3, 4, 8, 14) # BODY, OBJECT, CONTENT_TITLE_BODY
        )
        is_lorem_ipsum = "lorem ipsum" in shape.text.lower()
        is_empty_text_box = not shape.text.strip() and shape.height > Pt(100) # Heuristic for larger empty text boxes

        if not body_populated and (is_body_placeholder or is_lorem_ipsum or is_empty_text_box):
            tf = shape.text_frame
            tf.clear() # Clear existing content
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = content.get("body", "")
            # No longer forcing bold here.
            body_populated = True

        if title_populated and body_populated:
            break # Exit loop once both title and body content are placed

# --- Streamlit App ---
st.set_page_config(page_title="Dynamic AI Presentation Assembler", layout="wide")
st.title("📊 Dynamic AI Presentation Assembler") # Updated emoji for presentation

with st.sidebar:
    st.header("1. API Key & Decks")
    api_key = st.text_input("OpenAI API Key", type="password")
    st.markdown("---")
    st.header("2. Upload Decks")
    template_files = st.file_uploader("Upload Template Deck(s)", type=["pptx"], accept_multiple_files=True)
    gtm_file = st.file_uploader("Upload GTM Global Deck", type=["pptx"])
    st.markdown("---")
    st.header("3. Define Presentation Structure")
    
    # Initialize session state for structure if not present
    if 'structure' not in st.session_state: 
        st.session_state.structure = []
    
    # Button to add a new step to the presentation structure
    if st.button("Add New Step", use_container_width=True):
        st.session_state.structure.append({"id": str(uuid.uuid4()), "keyword": "", "action": "Copy from GTM (as is)"})

    # Display and manage each step in the structure
    for i, step in enumerate(st.session_state.structure):
        with st.container(border=True): # Use a container for visual separation
            cols = st.columns([3, 3, 1]) # Three columns for keyword, action, and delete button
            # Text input for the slide type keyword
            step["keyword"] = cols[0].text_input("Slide Type", step["keyword"], key=f"keyword_{step['id']}")
            # Selectbox for the action to perform (Copy or Merge)
            step["action"] = cols[1].selectbox(
                "Action", 
                ["Copy from GTM (as is)", "Merge: Template Layout + GTM Content"], 
                index=["Copy from GTM (as is)", "Merge: Template Layout + GTM Content"].index(step["action"]), 
                key=f"action_{step['id']}"
            )
            # Delete button for each step
            if cols[2].button("🗑️", key=f"del_{step['id']}"): # Changed emoji for delete
                st.session_state.structure.pop(i) # Remove the step
                st.rerun() # Rerun to update the UI immediately

    # Button to clear all defined steps
    if st.button("Clear Structure", use_container_width=True): 
        st.session_state.structure = []
        st.rerun()

# --- Main App Logic ---
# Check if all necessary inputs are provided before enabling assembly
if template_files and gtm_file and api_key and st.session_state.structure:
    # Button to trigger the presentation assembly process
    if st.button("🚀 Assemble Presentation", type="primary"): # Changed emoji for assemble
        with st.spinner("Assembling your new presentation..."):
            try:
                st.write("Step 1/3: Loading decks...")
                # CRITICAL: Use the first uploaded template file as the base for the new presentation.
                new_prs = Presentation(io.BytesIO(template_files[0].getvalue()))
                gtm_prs = Presentation(io.BytesIO(gtm_file.getvalue()))
                
                process_log = [] # To store logs of what happened during assembly
                st.write("Step 2/3: Building new presentation based on your structure...")
                
                num_template_slides = len(new_prs.slides)
                num_structure_steps = len(st.session_state.structure)

                # Prune excess slides from the template if the defined structure is shorter
                if num_structure_steps < num_template_slides:
                    # Iterate backwards to safely delete slides
                    for i in range(num_template_slides - 1, num_structure_steps - 1, -1):
                        rId = new_prs.slides._sldIdLst[i].rId # Get relationship ID
                        new_prs.part.drop_rel(rId) # Drop relationship
                        del new_prs.slides._sldIdLst[i] # Delete slide from slide list
                    st.info(f"Removed {num_template_slides - num_structure_steps} unused slides from the template.")
                elif num_structure_steps > num_template_slides:
                     st.warning(f"Warning: Your defined structure has more steps ({num_structure_steps}) than the template has slides ({num_template_slides}). Extra steps will be ignored.")

                # Process slides based on the defined structure
                for i, step in enumerate(st.session_state.structure):
                    # Ensure we don't go out of bounds if the template was trimmed or structure is longer
                    if i >= len(new_prs.slides): 
                        break

                    dest_slide = new_prs.slides[i] # Get the current destination slide from the new presentation
                    keyword, action = step["keyword"], step["action"]
                    log_entry = {"step": i + 1, "keyword": keyword, "action": action, "log": []}
                    
                    if action == "Copy from GTM (as is)":
                        # Find the best matching slide in the GTM deck using AI
                        result = find_slide_by_ai(api_key, gtm_prs, keyword, "GTM Deck")
                        log_entry["log"].append(f"**GTM Content Choice Justification:** {result['justification']}")
                        if result["slide"]:
                            # If a suitable slide is found, deep copy its content to the destination slide
                            deep_copy_slide_content(dest_slide, result["slide"])
                            log_entry["log"].append(f"**Action:** Replaced Template slide {i + 1} with content from GTM slide {result['index'] + 1}.")
                        else:
                            log_entry["log"].append("**Action:** No suitable slide found in GTM deck. Template slide was left as is.")
                    
                    elif action == "Merge: Template Layout + GTM Content":
                        # Find the best matching slide for content in the GTM deck using AI
                        content_result = find_slide_by_ai(api_key, gtm_prs, keyword, "GTM Deck")
                        log_entry["log"].append(f"**GTM Content Choice Justification:** {content_result['justification']}")
                        if content_result["slide"]:
                            # If content slide found, extract its title and body
                            content = get_slide_content(content_result["slide"])
                            # Populate the destination slide (template layout) with the extracted content
                            populate_slide(dest_slide, content)
                            log_entry["log"].append(f"**Action:** Merged content from GTM slide {content_result['index'] + 1} into Template slide {i+1}.")
                        else:
                             log_entry["log"].append("**Action:** No suitable content found in GTM deck. Template slide was left as is.")
                    
                    process_log.append(log_entry) # Add step log to overall process log
 
                st.success("Successfully built the new presentation structure.")
                
                st.write("Step 3/3: Finalizing...")
                st.subheader("📋 Process Log") # Changed emoji for process log
                # Display the process log in an expandable format
                for entry in process_log:
                    with st.expander(f"Step {entry['step']}: '{entry['keyword']}' ({entry['action']})"):
                        for line in entry['log']: 
                            st.markdown(f"- {line}")
                
                # Save the assembled presentation to an in-memory buffer
                output_buffer = io.BytesIO()
                new_prs.save(output_buffer)
                output_buffer.seek(0) # Rewind the buffer to the beginning for downloading

                st.success("✨ Your new regional presentation has been assembled!") # Changed emoji for success
                # Provide a download button for the user
                st.download_button(
                    "Download Assembled PowerPoint", 
                    data=output_buffer, 
                    file_name="Dynamic_AI_Assembled_Deck.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"A critical error occurred: {e}")
                st.exception(e) # Display full traceback for debugging
else:
    # Instructions displayed when inputs are not yet complete
    st.info("Please provide an API Key, upload at least one Template Deck and a GTM Deck, and define the structure in the sidebar to begin.")
