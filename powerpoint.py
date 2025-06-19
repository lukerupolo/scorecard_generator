from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import matplotlib.pyplot as plt
import requests 
import streamlit as st
import pandas as pd
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ================================================================================
# Main Presentation Creation Function
# ================================================================================
def create_presentation(title, subtitle, scorecard_moments, sheets_dict, style_guide, region_prompt, openai_api_key):
    """Creates and returns a PowerPoint presentation as a BytesIO buffer."""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    add_title_slide(prs, title, subtitle, style_guide, region_prompt, openai_api_key)
    add_timeline_slide(prs, scorecard_moments, style_guide)

    total_moments = len(scorecard_moments)
    if total_moments > 0:
        progress_text = "Generating AI background images... (This can take a moment)"
        image_progress_bar = st.progress(0, text=progress_text)

        for i, moment in enumerate(scorecard_moments):
            image_progress_bar.progress((i + 1) / total_moments, text=f"Generating image for '{moment}'...")
            add_moment_title_slide(prs, f"SCORECARD:\n{moment.upper()}", style_guide, region_prompt, openai_api_key)
            for sheet_name, scorecard_df in sheets_dict.items():
                if "benchmark" not in sheet_name.lower():
                    add_df_to_slide(prs, scorecard_df, f"{moment.upper()} METRICS: {sheet_name}", style_guide)
        
        image_progress_bar.empty()

    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# ================================================================================
# AI Background Image Generation
# ================================================================================
def generate_and_add_background_image(slide, region, style_guide, api_key, slide_width, slide_height, prompt_detail="football culture"):
    prompt = f"Dark, gritty, artistic representation of {prompt_detail} in {region}, cinematic, ultra-realistic photo, dramatic lighting, epic style"
    if not api_key:
        st.warning("OpenAI API key is missing. Using a solid background.")
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = style_guide["colors"]["title_slide_bg"]
        return
    try:
        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
        payload = {"model": "dall-e-3", "prompt": prompt, "n": 1, "size": "1792x1024", "response_format": "url"}
        api_url = "https://api.openai.com/v1/images/generations"
        response = requests.post(api_url, headers=headers, json=payload, timeout=45)
        response.raise_for_status()
        image_url = response.json()['data'][0]['url']
        image_response = requests.get(image_url, timeout=15); image_response.raise_for_status()
        image_stream = BytesIO(image_response.content)
        pic = slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=slide_width, height=slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
    except requests.exceptions.RequestException as e:
        st.error(f"Image generation for '{region}' failed: {e}. Using a solid background.")
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = style_guide["colors"]["title_slide_bg"]

# ================================================================================
# Helper functions for slide creation and styling
# ================================================================================
def add_title_slide(prs, title_text, subtitle_text, style_guide, region, api_key):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    generate_and_add_background_image(slide, region, style_guide, api_key, prs.slide_width, prs.slide_height, prompt_detail="a cinematic football stadium")
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
    p = title_shape.text_frame.paragraphs[0]; p.text = title_text.upper(); p.font.name = style_guide["fonts"]["heading"]; p.font.bold = True; p.font.size = style_guide["font_sizes"]["title"]; p.font.color.rgb = style_guide["colors"]["title_slide_text"]; p.alignment = PP_ALIGN.CENTER
    subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1.5))
    p = subtitle_shape.text_frame.paragraphs[0]; p.text = subtitle_text; p.font.name = style_guide["fonts"]["body"]; p.font.size = style_guide["font_sizes"]["subtitle"]; p.font.color.rgb = style_guide["colors"]["title_slide_text"]; p.alignment = PP_ALIGN.CENTER

def add_moment_title_slide(prs, title_text, style_guide, region, api_key):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    generate_and_add_background_image(slide, region, style_guide, api_key, prs.slide_width, prs.slide_height)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(3))
    p = txBox.text_frame.paragraphs[0]; p.text = title_text; p.font.name = style_guide["fonts"]["heading"]; p.font.bold = True; p.font.size = style_guide["font_sizes"]["moment_title"]; p.font.color.rgb = style_guide["colors"]["title_slide_text"]; p.alignment = PP_ALIGN.CENTER

def add_timeline_slide(prs, timeline_moments, style_guide):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = style_guide["colors"]["content_slide_bg"]
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
    p = title_shape.text_frame.paragraphs[0]; p.text = "TIMELINE"; p.font.name = style_guide["fonts"]["heading"]; p.font.bold = True; p.font.size = style_guide["font_sizes"]["title"]; p.font.color.rgb = style_guide["colors"]["content_heading_text"]; p.alignment = PP_ALIGN.CENTER
    if not timeline_moments: return
    fig, ax = plt.subplots(figsize=(14, 2.5))
    fig.patch.set_facecolor(f'#{style_guide["colors"]["content_slide_bg"]}')
    ax.set_facecolor(f'#{style_guide["colors"]["content_slide_bg"]}')
    ax.axhline(0, color=f'#{style_guide["colors"]["content_body_text"]}', xmin=0.05, xmax=0.95, zorder=1, linewidth=1.5)
    for i, moment in enumerate(timeline_moments):
        ax.plot(i + 1, 0, 'o', markersize=20, color=f'#{style_guide["colors"]["content_heading_text"]}', zorder=2)
        ax.text(x=i + 1, y=-0.3, s=moment.upper(), ha='center', va='top', fontsize=12, fontname='sans-serif', color=f'#{style_guide["colors"]["content_body_text"]}', weight='bold')
    ax.set_ylim(-1, 1); ax.axis('off'); plt.tight_layout(pad=0.1)
    plot_stream = BytesIO(); plt.savefig(plot_stream, format='png', facecolor=fig.get_facecolor(), transparent=False); plt.close(fig); plot_stream.seek(0)
    slide.shapes.add_picture(plot_stream, Inches(1), Inches(3.5), width=Inches(14))

def apply_table_style_pptx(table, style_guide):
    """
    Styles a table in PowerPoint using the provided style guide.
    """
    header_bg = style_guide["colors"]["table_header_bg"]
    header_text = style_guide["colors"]["table_header_text"]
    body_text = style_guide["colors"]["content_body_text"]
    row_bg = style_guide["colors"]["table_alt_row_bg"] # Use this key for the row background
    
    heading_font, body_font = style_guide["fonts"]["heading"], style_guide["fonts"]["body"]
    header_fs, body_fs = style_guide["font_sizes"]["table_header"], style_guide["font_sizes"]["table_body"]
    
    # Style header row
    for cell in table.rows[0].cells:
        cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]; p.font.color.rgb = header_text; p.font.name = heading_font; p.font.size = header_fs; p.alignment = PP_ALIGN.CENTER
    
    # Style data rows
    for i, row in enumerate(table.rows):
        if i == 0: continue # Skip header row
        
        # FIXED: Apply the same background color to ALL data rows
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            
            # Apply text styling
            p = cell.text_frame.paragraphs[0]
            p.font.name = body_font
            p.font.size = body_fs
            p.font.color.rgb = body_text

def add_df_to_slide(prs, df, slide_title, style_guide):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = style_guide["colors"]["content_slide_bg"]
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
    p = title_shape.text_frame.paragraphs[0]; p.text = slide_title; p.font.name = style_guide['fonts']['heading']; p.font.size = style_guide['font_sizes']['content_title']; p.font.color.rgb = style_guide['colors'].get("content_heading_text")

    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.2), Inches(15), Inches(1.0)).table
    table.columns[0].width = Inches(2.0); table.columns[1].width = Inches(4.5)
    for i in range(2, cols): table.columns[i].width = Inches(2.0)

    table.cell(0, 0).text = ""
    for i, col_name in enumerate(df.columns[1:], start=1): table.cell(0, i).text = col_name

    for r in range(rows):
        for c in range(cols): table.cell(r + 1, c).text = str(df.iloc[r, c])
    
    apply_table_style_pptx(table, style_guide)

    df_copy = df.copy()
    if 'Category' in df_copy.columns:
        df_copy['category_group'] = (df_copy['Category'] != '').cumsum()
        for group_id in df_copy['category_group'].unique():
            group_rows = df_copy[df_copy['category_group'] == group_id]
            if len(group_rows) > 1:
                start_row_idx = group_rows.index[0] + 1; end_row_idx = group_rows.index[-1] + 1
                start_cell = table.cell(start_row_idx, 0); end_cell = table.cell(end_row_idx, 0)
                start_cell.merge(end_cell)
        
        for r in range(1, rows + 1):
            cell = table.cell(r, 0)
            if cell.text:
                p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(14); p.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
